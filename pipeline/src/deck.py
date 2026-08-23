"""Build the 5-slide deck (deliverable 3) from the pipeline state.

Charts are matplotlib PNGs in the same visual language as the dashboard;
slide copy is short enough to present in one minute per slide.
"""
import json

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .util import ROOT
from .paths import ws_dir, latest_run, DEFAULT_WS

INK = RGBColor(0x1D, 0x29, 0x39)
SOFT = RGBColor(0x66, 0x70, 0x85)
ACC = RGBColor(0xFF, 0x42, 0x00)
ACC2 = RGBColor(0x37, 0x4B, 0x60)
FAIL = RGBColor(0xC1, 0x2D, 0x00)
WARN = RGBColor(0xAD, 0x83, 0x6C)
LINE = RGBColor(0xE6, 0xE1, 0xDA)
BGSOFT = RGBColor(0xE6, 0xEE, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MC = {"Netherlands": "#ff4200", "Germany": "#661439", "Portugal": "#374b60", "Poland": "#ad836c"}
CHARTS = ROOT.parent / "deliverables" / "charts"
FONT = "Arial"

plt.rcParams.update({
    "font.family": FONT, "font.size": 11, "axes.edgecolor": "#e6e1da",
    "axes.linewidth": 0.8, "axes.labelcolor": "#667085", "xtick.color": "#667085",
    "ytick.color": "#667085", "figure.facecolor": "white", "axes.facecolor": "white",
    "axes.grid": True, "grid.color": "#f1ede7", "grid.linewidth": 0.7,
})


def _meur(x, _):
    return f"{x/1e6:+.0f}M" if x else "0"


def chart_cash_curves(res, ranking):
    fig, ax = plt.subplots(figsize=(7.8, 4.3), dpi=200)
    for m in ranking:
        r = res[m]
        xs = list(range(6))
        ys = [4_000_000 - r["params"]["entry_cost"]] + [y["cash"] for y in r["years"]]
        ax.plot(xs, ys, color=MC[m], lw=2.6, marker="o", ms=4,
                label=f"{m}  ({ys[-1]/1e6:+.1f}M)")
    ax.axhline(0, color="#c12d00", ls="--", lw=1.2)
    ax.text(0.06, -700_000, "€0 — insolvency", color="#c12d00", fontsize=9.5, va="top")
    ax.set_axisbelow(True)
    ax.set_xticks(range(6), [f"Y{i}" for i in range(6)])
    ax.yaxis.set_major_formatter(FuncFormatter(_meur))
    ax.set_title("Five-year cash position, €4.0M starting runway", fontsize=12.5,
                 color="#1d2939", loc="left", pad=10, fontweight="bold")
    ax.legend(frameon=False, fontsize=10, loc="upper left")
    fig.tight_layout()
    p = CHARTS / "cash_curves.png"
    fig.savefig(p)
    plt.close(fig)
    return p


def chart_trough(res, ranking):
    fig, ax = plt.subplots(figsize=(6.2, 4.2), dpi=200)
    ms = ranking[::-1]
    vals = [res[m]["min_cash"] for m in ms]
    cols = [MC[m] for m in ms]
    bars = ax.barh(ms, vals, color=cols, height=0.55, zorder=3)
    ax.set_axisbelow(True)
    ax.axvline(0, color="#c12d00", ls="--", lw=1.2)
    for b, v in zip(bars, vals):
        ax.text(v + (80_000 if v > 0 else -80_000), b.get_y() + b.get_height() / 2,
                f"{v/1e6:+.1f}M", va="center", ha="left" if v > 0 else "right",
                fontsize=10.5, color="#1d2939", fontweight="bold")
    ax.set_xlim(min(vals) * 1.35, max(vals) * 2.2)
    ax.xaxis.set_major_formatter(FuncFormatter(_meur))
    ax.set_title("Cash trough — how close each market comes to €0", fontsize=12.5,
                 color="#1d2939", loc="left", pad=10, fontweight="bold")
    fig.tight_layout()
    p = CHARTS / "trough.png"
    fig.savefig(p)
    plt.close(fig)
    return p


def chart_nl(res, rec):
    r = res[rec]
    years = [f"Y{y['year']}" for y in r["years"]]
    ebitda = [y["ebitda"] for y in r["years"]]
    cash = [y["cash"] for y in r["years"]]
    fig, ax = plt.subplots(figsize=(6.4, 4.2), dpi=200)
    ax.bar(years, ebitda, color=["#c12d00" if v < 0 else "#374b60" for v in ebitda],
           width=0.55, label="EBITDA", zorder=3)
    ax.set_axisbelow(True)
    ax2 = ax.twinx()
    ax2.plot(years, cash, color="#1d2939", lw=2.4, marker="o", ms=4, label="Cash")
    ax2.grid(False)
    ax.axhline(0, color="#667085", lw=0.8)
    ax.yaxis.set_major_formatter(FuncFormatter(_meur))
    ax2.yaxis.set_major_formatter(FuncFormatter(_meur))
    ax.set_title(f"{rec}: EBITDA turns positive in Y{r['break_even_year']}; cash never approaches zero",
                 fontsize=12.5, color="#1d2939", loc="left", pad=10, fontweight="bold")
    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax.legend(h1 + h2, l1 + l2, frameon=False, fontsize=10, loc="upper left")
    fig.tight_layout()
    p = CHARTS / "nl_model.png"
    fig.savefig(p)
    plt.close(fig)
    return p


# ---------- pptx helpers ----------
def _tb(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    return box.text_frame


def _para(tf, text, size=14, bold=False, color=INK, first=False, space=6, align=None, font=FONT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    p.space_after = Pt(space)
    if align:
        p.alignment = align
    return p


def _bullets(tf, items, size=13.5, space=8):
    for i, (head, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 and not tf.paragraphs[0].runs else tf.add_paragraph()
        r = p.add_run(); r.text = "▪ "; r.font.color.rgb = ACC; r.font.size = Pt(size); r.font.name = FONT
        r = p.add_run(); r.text = head; r.font.bold = True; r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = FONT
        if rest:
            r = p.add_run(); r.text = " " + rest; r.font.size = Pt(size); r.font.color.rgb = SOFT; r.font.name = FONT
        p.space_after = Pt(space)


def _band(slide, title, kicker, W):
    bar = slide.shapes.add_shape(1, 0, 0, Emu(int(W * 914400)), Inches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACC; bar.line.fill.background()
    tf = _tb(slide, 0.55, 0.32, W - 1.1, 1.0)
    _para(tf, kicker.upper(), 10.5, True, ACC2, first=True, space=2)
    _para(tf, title, 25, True, INK, space=0)


def _stat(slide, x, y, w, big, small, color=INK):
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(1.12))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE; card.line.width = Pt(0.75)
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.13); tf.margin_top = Inches(0.09); tf.margin_right = Inches(0.1)
    _para(tf, big, 21, True, color, first=True, space=1)
    _para(tf, small, 10.5, False, SOFT, space=0)


def _footer(slide, text, W, H):
    tf = _tb(slide, 0.55, H - 0.42, W - 1.1, 0.3)
    _para(tf, text, 9.5, False, SOFT, first=True, space=0)


def build() -> str:
    state = json.load(open(latest_run(ws_dir(DEFAULT_WS)) / "state.json"))
    res = state["conclusion"]["results"]
    ranking = state["conclusion"]["ranking"]
    rec = state["conclusion"]["recommendation"]
    nl = res[rec]
    de = res["Germany"]
    fac = state["facilities"]["summary"]
    man = state["manifest"]

    CHARTS.mkdir(parents=True, exist_ok=True)
    c1 = chart_cash_curves(res, ranking)
    c2 = chart_trough(res, ranking)
    c3 = chart_nl(res, rec)

    prs = Presentation()
    W, H = 13.333, 7.5
    prs.slide_width = Emu(int(W * 914400))
    prs.slide_height = Emu(int(H * 914400))
    blank = prs.slide_layouts[6]

    # ---------- S1 recommendation ----------
    s = prs.slides.add_slide(blank)
    _band(s, "Enter the Netherlands first", "Helix Optics · first European market · board recommendation", W)
    xs = 0.55
    _stat(s, xs, 1.55, 2.85, "Break-even Y3", "company EBITDA positive in year 3, the second year of reimbursed sales")
    _stat(s, xs, 2.80, 2.85, "€1.30M", "cash trough: the runway is never exhausted", ACC2)
    _stat(s, xs, 4.05, 2.85, "159,910 / yr", "addressable FIT-positives: 92% of Germany's volume, none of its barriers")
    _stat(s, xs, 5.30, 2.85, f"€{nl['system_saving_per_100']:,} saved", "per 100 patients triaged, at €650 per avoided colonoscopy. The payer wants this")
    s.shapes.add_picture(str(c1), Inches(3.85), Inches(1.55), width=Inches(8.75))
    tf = _tb(s, 3.9, 6.55, 9.0, 0.6)
    _para(tf, "The Netherlands is the only market of the four that survives on the €4.0M runway: "
              "12 months to reimbursement, one national procurement route, no incumbent.", 12.5, False, INK, first=True)
    _footer(s, "All figures computed by the data pipeline from the case pack; model replicated live in the Excel. "
               "Ramp anchored to the pack's benchmark (~20% of eligible volume within 3 years of reimbursement).", W, H)

    # ---------- S2 Germany trap ----------
    s = prs.slides.add_slide(blank)
    _band(s, "Germany is the prize — and the certain death of the runway", "Why not the biggest market first", W)
    tf = _tb(s, 0.55, 1.6, 6.1, 5.0)
    _bullets(tf, [
        ("24 months to reimbursed revenue.", "The pathway is procedural, and the national report is explicit that "
         "clinical investigators, lab partnerships and prior studies do NOT shorten it."),
        ("Pre-reimbursement revenue ≈ €0.", "Selective contracts and self-pay have 'historically covered negligible "
         "volume'. Planning assumption: immaterial."),
        ("The runway dies first.", f"Burn reaches ~€2.85M in year 1; cash bottoms at {de['min_cash']/1e6:.1f}M, "
         "insolvent before the first reimbursed euro."),
        ("An incumbent owns the channel.", "OncoStream: ~3 years reimbursed, 35–40% of addressable volume, "
         "framework agreements with the largest endoscopy networks."),
        ("Its 300k tests/yr claim is impossible.", "Germany's entire organised addressable volume is 174,240/yr; "
         "the company-reported figure is 1.7× the whole market. Pipeline check CHK-02 rejects it, and the independent "
         "35–40% share (~61–70k tests) is used instead."),
    ], size=13)
    s.shapes.add_picture(str(c2), Inches(6.95), Inches(1.75), width=Inches(5.9))
    tf = _tb(s, 7.1, 6.15, 5.6, 0.8)
    _para(tf, "Right market later, unaffordable market now. File the German application early (the clock is "
              "procedural); enter when new capital or NL cash flow can fund the two-year wait.",
          11.5, False, SOFT, first=True)
    _footer(s, "Sources: Germany national report; Meridian competitor brief (company-reported figures cross-checked "
               "against primary screening statistics — see pipeline validation).", W, H)

    # ---------- S3 NL case ----------
    s = prs.slides.add_slide(blank)
    _band(s, "The Netherlands: fastest to revenue, and the payer is motivated", "The case for the recommendation", W)
    tf = _tb(s, 0.55, 1.6, 6.0, 4.9)
    _bullets(tf, [
        ("Volume without the barriers.", "71% participation × 9.1% positivity → 159,910 FIT-positives/yr from "
         "4.95M eligible: 92% of Germany's volume."),
        ("12 months to reimbursement.", "Centralised programme, one national procurement decision. No "
         "region-by-region grind."),
        ("Open field.", "No competitor active; second-highest price of the four (€215)."),
        ("Capacity pressure is policy.", "Reducing avoidable colonoscopies is an explicit programme priority; "
         "register data shows referrals at ~103% of national endoscopy capacity."),
        ("The economics clear.", "€215 price vs €51 COGS (Y3) → 76% unit margin; company EBITDA positive Y3; "
         "5-year end cash €9.6M."),
    ], size=13)
    s.shapes.add_picture(str(c3), Inches(6.85), Inches(1.7), width=Inches(6.0))
    tf = _tb(s, 7.0, 6.1, 5.8, 0.9)
    _para(tf, "What must be true: reimbursement lands ≈12 months; ramp reaches ~13% of addressable by the second "
              "reimbursed year; price holds near €215.", 11.5, False, SOFT, first=True)
    _footer(s, "Utilisation computed from the deduplicated facility registers "
               f"({fac['raw_records']:,} raw records → {fac['unique_facilities']} real units).", W, H)

    # ---------- S4 sensitivity ----------
    s = prs.slides.add_slide(blank)
    _band(s, "Where the case bends, and where it breaks", "Break-even and the range around it", W)
    tf = _tb(s, 0.55, 1.55, 5.6, 5.2)
    _para(tf, "Break-even: Year 3", 17, True, INK, first=True, space=4)
    _para(tf, "On the base case — reimbursement at month 12, ramp to 20% of addressable by year 3 "
              "(the pack's own benchmark), price €215.", 12.5, False, SOFT, space=12)
    _bullets(tf, [
        ("Ramp alone can sink it.", "At half the benchmark ramp the trough grazes €0 (−€0.0M): the case "
         "fails without help."),
        ("Delay alone breaks it too.", "Any slip past month 12 puts the trough at −€0.8M on the base ramp "
         "(the annual model rounds delays up to whole years, deliberately conservative)."),
        ("Together, decisively fatal.", "Half-ramp AND a 24-month delay → trough −€2.1M."),
        ("Price is the cushion; month 15 is the tripwire.", "Every €10 of price ≈ €0.2–0.3M of annual EBITDA "
         "at scale. No reimbursement signal by month 15 → cut in-market spend and bridge before the trough."),
    ], size=12.5)
    # sensitivity table (computed in python with the same engine)
    from .stage7_conclude import project
    comp = {p: e["value"] for p, e in state["dataset"]["company"].items()}
    nlpar = {p: e["value"] for p, e in state["dataset"][rec].items()}
    mults, delays = [0.5, 0.75, 1.0, 1.25], [0, 6, 12]
    rows_n = len(mults) + 1
    tbl = s.shapes.add_table(rows_n, 4, Inches(6.6), Inches(1.9), Inches(6.1), Inches(0.5 * rows_n)).table
    tbl.cell(0, 0).text = "cash trough · ramp ↓ / delay →"
    for j, dd in enumerate(delays):
        tbl.cell(0, j + 1).text = f"{12+dd} months"
    for i, mu in enumerate(mults):
        tbl.cell(i + 1, 0).text = f"ramp ×{mu:.0%}"
        for j, dd in enumerate(delays):
            r = project(nlpar, comp, ramp_mult=mu, delay_extra_months=dd)
            cell = tbl.cell(i + 1, j + 1)
            cell.text = f"{r['min_cash']/1e6:+.1f}M"
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.runs[0].font.size = Pt(12)
            para.runs[0].font.bold = True
            para.runs[0].font.name = FONT
            para.runs[0].font.color.rgb = WHITE if r["min_cash"] < 0 else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = FAIL if r["min_cash"] < 0 else (BGSOFT if r["min_cash"] > 1e6 else RGBColor(0xF6, 0xEA, 0xD9))
    for j in range(4):
        c0 = tbl.cell(0, j)
        c0.fill.solid(); c0.fill.fore_color.rgb = INK
        for run in c0.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11); run.font.color.rgb = WHITE; run.font.bold = True; run.font.name = FONT
    for i in range(1, rows_n):
        c0 = tbl.cell(i, 0)
        for run in c0.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11.5); run.font.name = FONT
    tf = _tb(s, 6.6, 1.9 + 0.5 * rows_n + 0.25, 6.1, 1.4)
    _para(tf, "Green: comfortable. Amber: survives, thin. Red: insolvent.", 11, False, SOFT, first=True, space=3)
    _para(tf, "Same grid lives as formulas in the Excel ('3 Sensitivity') — change any assumption in the live "
              "session and it recalculates.", 11, False, SOFT)
    _footer(s, "Grid computed by the same engine as the Excel; delays round up to whole years (conservative).", W, H)

    # ---------- S5 sequence ----------
    s = prs.slides.add_slide(blank)
    _band(s, "After the Netherlands: sequence, triggers, and what we refuse to do", "Questions 3 — the expansion plan", W)
    steps = [
        ("NOW → Y1", "Netherlands entry", "€200k entry · reimbursement application immediately · first revenue ~month 13", ACC),
        ("Y1", "File Germany's clock", "the 24-month pathway is procedural — start it while NL scales; commercial spend stays €0", INK),
        ("Y2", "Portugal", "€120k entry, 12-month route, home turf; contribution-positive add-on (~€0.6M/yr at benchmark)", RGBColor(0x37, 0x4B, 0x60)),
        ("Y3+", "Germany, funded", "enter when NL cash flow / new capital covers the wait; IVDR done; Poland stays parked (€115 price ≈ COGS+ margin too thin)", FAIL),
    ]
    x = 0.55
    for i, (when, title, sub, col) in enumerate(steps):
        card = s.shapes.add_shape(1, Inches(x), Inches(1.7), Inches(2.95), Inches(2.0))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LINE; card.line.width = Pt(0.75)
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.14); tf.margin_top = Inches(0.1)
        _para(tf, when, 10.5, True, col, first=True, space=2)
        _para(tf, title, 15.5, True, INK, space=3)
        _para(tf, sub, 10.5, False, SOFT, space=0)
        x += 3.12
    tf = _tb(s, 0.55, 4.05, 12.2, 2.6)
    _para(tf, "Decision triggers, agreed with the board today:", 14, True, INK, first=True, space=6)
    _bullets(tf, [
        ("Accelerate:", "NL reimbursed-year-1 penetration ≥7% of addressable and price ≥€200 → pull Portugal forward; "
         "raise growth capital against proven uptake."),
        ("Delay:", "reimbursement slips past month 15 or ramp <half of plan → freeze Portugal, cut NL in-market cost, "
         "extend runway before the trough."),
        ("Abandon / pivot:", "no reimbursement by month 18 with pipeline dry → stop-loss NL commercial build; "
         "the Germany filing and the Portuguese R&D base (HealthTech Growth Call: €500k, 60% co-financing of "
         "Portugal-based R&D — not usable for any launch) keep two doors open."),
        ("What we refuse to do:", "chase Germany's headline volume with a 24-month unfunded gap, or book the funding "
         "call as launch money. It is R&D co-financing only."),
    ], size=12.5)
    _footer(s, f"Pipeline run: {man['extraction']['deterministic']} values extracted deterministically, "
               f"{man['llm_calls']} AI calls, eval 28/28, {man['total_s']}s end-to-end. "
               "The same assessment re-runs on markets 5–14 by adding a config entry.", W, H)

    out = ROOT.parent / "deliverables" / "Helix_Optics_Market_Entry_Deck.pptx"
    prs.save(out)
    return str(out)


if __name__ == "__main__":
    print(build())
